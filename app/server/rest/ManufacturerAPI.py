from flask import Flask, jsonify, abort, make_response, request
from flask.ext.restful import Api, Resource, reqparse, fields, marshal_with
from server.dao import ManufacturerDAO
import json
from server import app, api
import cairosvg
from pptx import Presentation
from pptx.util import Inches

resource_fields = {
'pkey': fields.Integer,
'name': fields.String
}

print(__name__)

class ManufacturerAPI(Resource):
    @marshal_with(resource_fields)
    def get(self):
        manufacturers = ManufacturerDAO.getManufacturers()
        return manufacturers

class TestAPI(Resource):
    def get(self):
        svgPNG = '<svg xmlns="http://www.w3.org/2000/svg" id="chart" class="chart" width="420" height="120"><g transform="translate(0,0)"><rect width="40" height="19" fill="steelblue"/><text x="37" y="10" dy=".35em">4</text></g><g transform="translate(0,20)"><rect width="80" height="19" fill="steelblue"/><text x="77" y="10" dy=".35em">8</text></g><g transform="translate(0,40)"><rect width="150" height="19" fill="steelblue"/><text x="147" y="10" dy=".35em">15</text></g><g transform="translate(0,60)"><rect width="160" height="19" fill="steelblue"/><text x="157" y="10" dy=".35em">16</text></g><g transform="translate(0,80)"><rect width="230.00000000000003" height="19" fill="steelblue"/><text x="227.00000000000003" y="10" dy=".35em">23</text></g><g transform="translate(0,100)"><rect width="420" height="19" fill="steelblue"/><text x="417" y="10" dy=".35em">42</text></g></svg>'
        svgPDF = '<svg xmlns="http://www.w3.org/2000/svg"><svg xmlns="http://www.w3.org/2000/svg" id="chart" class="chart" width="420" height="120"><g transform="translate(0,0)"><rect width="40" height="19" fill="steelblue"/><text x="37" y="10" dy=".35em">4</text></g><g transform="translate(0,20)"><rect width="80" height="19" fill="steelblue"/><text x="77" y="10" dy=".35em">8</text></g><g transform="translate(0,40)"><rect width="150" height="19" fill="steelblue"/><text x="147" y="10" dy=".35em">15</text></g><g transform="translate(0,60)"><rect width="160" height="19" fill="steelblue"/><text x="157" y="10" dy=".35em">16</text></g><g transform="translate(0,80)"><rect width="230.00000000000003" height="19" fill="steelblue"/><text x="227.00000000000003" y="10" dy=".35em">23</text></g><g transform="translate(0,100)"><rect width="420" height="19" fill="steelblue"/><text x="417" y="10" dy=".35em">42</text></g></svg><svg xmlns="http://www.w3.org/2000/svg" id="chart" class="chart" width="420" height="120"><g transform="translate(0,0)"><rect width="40" height="19" fill="steelblue"/><text x="37" y="10" dy=".35em">4</text></g><g transform="translate(0,20)"><rect width="80" height="19" fill="steelblue"/><text x="77" y="10" dy=".35em">8</text></g><g transform="translate(0,40)"><rect width="150" height="19" fill="steelblue"/><text x="147" y="10" dy=".35em">15</text></g><g transform="translate(0,60)"><rect width="160" height="19" fill="steelblue"/><text x="157" y="10" dy=".35em">16</text></g><g transform="translate(0,80)"><rect width="230.00000000000003" height="19" fill="steelblue"/><text x="227.00000000000003" y="10" dy=".35em">23</text></g><g transform="translate(0,100)"><rect width="420" height="19" fill="steelblue"/><text x="417" y="10" dy=".35em">42</text></g></svg><svg xmlns="http://www.w3.org/2000/svg" id="chart" class="chart" width="420" height="120"><g transform="translate(0,0)"><rect width="40" height="19" fill="steelblue"/><text x="37" y="10" dy=".35em">4</text></g><g transform="translate(0,20)"><rect width="80" height="19" fill="steelblue"/><text x="77" y="10" dy=".35em">8</text></g><g transform="translate(0,40)"><rect width="150" height="19" fill="steelblue"/><text x="147" y="10" dy=".35em">15</text></g><g transform="translate(0,60)"><rect width="160" height="19" fill="steelblue"/><text x="157" y="10" dy=".35em">16</text></g><g transform="translate(0,80)"><rect width="230.00000000000003" height="19" fill="steelblue"/><text x="227.00000000000003" y="10" dy=".35em">23</text></g><g transform="translate(0,100)"><rect width="420" height="19" fill="steelblue"/><text x="417" y="10" dy=".35em">42</text></g></svg></svg>'
        png = cairosvg.svg2png(bytestring=svgPNG)
        pdf = cairosvg.svg2pdf(bytestring=svgPDF)
        response = make_response(pdf)
        response.headers["Content-Disposition"] = "attachment; filename=test.pdf"
        return response

class PptxAPI(Resource):
    def get(self):
        # img_path = 'server/static/test.svg'
        img_path = 'server/static/login_screen.png'
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        left = top = Inches(1)
        pic = slide.shapes.add_picture(img_path, left, top)

        slide = prs.slides.add_slide(blank_slide_layout)

        left = top = Inches(1)
        pic = slide.shapes.add_picture(img_path, left, top)


        prs.save('test.pptx')
        ppt = open('test.pptx', 'rb')
        data = ppt.read() # if you only wanted to read 512 bytes, do .read(512)
        ppt.close()

        response = make_response(data)
        response.headers["Content-Disposition"] = "attachment; filename=ppt.pptx"
        return response





api.add_resource(ManufacturerAPI, '/api/v1.0/manufacturer', endpoint='manufacturer')
api.add_resource(TestAPI, '/test', endpoint='/test')
api.add_resource(PptxAPI, '/pptx', endpoint='/pptx')
